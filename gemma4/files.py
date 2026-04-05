import os
import csv
import pypdf
import docx
import openpyxl
import pptx
from typing import Optional
from gemma4.manager import get_manager

def read_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def read_pdf(file_path: str) -> str:
    text = ""
    with open(file_path, "rb") as f:
        reader = pypdf.PdfReader(f)
        for page in reader.pages:
            content = page.extract_text()
            if content:
                text += content + "\n"
    return text

def read_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def read_pptx(file_path: str) -> str:
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_xlsx(file_path: str) -> str:
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text = ""
    for sheet in wb.sheetnames:
        text += f"Sheet: {sheet}\n"
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
    return text

def read_csv(file_path: str) -> str:
    text = ""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            text += "\t".join(row) + "\n"
    return text

def read_file_content(file_path: str) -> str:
    """
    Tự động nhận diện định dạng và trích xuất nội dung văn bản từ file.
    """
    if not os.path.exists(file_path):
        return f"Lỗi: Không tìm thấy file tại {file_path}"
    
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == ".txt":
            return read_txt(file_path)
        elif ext == ".pdf":
            return read_pdf(file_path)
        elif ext in [".docx", ".doc"]:
            return read_docx(file_path)
        elif ext in [".pptx", ".ppt"]:
            return read_pptx(file_path)
        elif ext in [".xlsx", ".xls"]:
            return read_xlsx(file_path)
        elif ext == ".csv":
            return read_csv(file_path)
        else:
            # Fallback cho định dạng khác (thử đọc như text)
            return read_txt(file_path)
    except Exception as e:
        return f"Lỗi khi đọc file {ext}: {str(e)}"

def process_file_with_prompt(file_path: str, prompt: str, model_id: str = "google/gemma-4-e4b-it") -> str:
    """
    Đọc nội dung file và đưa vào Gemma 4 kèm theo prompt để xử lý.
    """
    content = read_file_content(file_path)
    if content.startswith("Lỗi:"):
        return content
    
    # Truncate content to avoid exceeding context limits (approx 4k tokens)
    # 1 token ~ 4 characters for English, ~2.5 characters for Vietnamese. 
    # Use 12000 chars as a conservative limit.
    max_chars = 12000
    if len(content) > max_chars:
        content = content[:max_chars] + "\n... [Nội dung đã được rút gọn do quá dài] ..."
        
    full_prompt = f"Nội dung của file:\n{content}\n\nYêu cầu: {prompt}"
    
    manager = get_manager(model_id)
    return manager.generate(full_prompt)
