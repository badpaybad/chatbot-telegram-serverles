import os
import sys
import re
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt

def download_image(url):
    """Tải ảnh từ URL và trả về BytesIO."""
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        return BytesIO(response.content)
    except Exception as e:
        print(f"Lỗi khi tải ảnh từ {url}: {e}")
        return None

class MarkdownToPPTConverter:
    def __init__(self, title="Presentation"):
        self.prs = Presentation()
        self.title = title
        self.current_slide = None
        self.current_content_placeholder = None

    def add_title_slide(self, title_text, subtitle_text=""):
        slide_layout = self.prs.slide_layouts[0]  # Layout 0 is Title Slide
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_text
        subtitle.text = subtitle_text
        self.current_slide = slide
        return slide

    def add_content_slide(self, title_text):
        slide_layout = self.prs.slide_layouts[1]  # Layout 1 is Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        self.current_slide = slide
        self.current_content_placeholder = slide.placeholders[1]
        return slide

    def add_bullet(self, text, level=0):
        if not self.current_content_placeholder:
            self.add_content_slide("Tiếp tục")
        
        p = self.current_content_placeholder.text_frame.add_paragraph()
        p.text = text
        p.level = level

    def add_image(self, path_or_url):
        if not self.current_slide:
            self.add_content_slide("Hình ảnh")
        
        image_data = None
        if path_or_url.startswith(('http://', 'https://')):
            image_data = download_image(path_or_url)
        elif os.path.exists(path_or_url):
            image_data = path_or_url
        
        if image_data:
            # Add image (heuristic position, center-ish)
            try:
                # Default position for simple demo
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                self.current_slide.shapes.add_picture(image_data, left, top, width=width)
            except Exception as e:
                print(f"Lỗi khi chèn ảnh: {e}")

    def add_video(self, path):
        if not self.current_slide:
            self.add_content_slide("Video")
        
        if os.path.exists(path):
            try:
                # Video embedding
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(4.5)
                # python-pptx add_movie requires a path
                self.current_slide.shapes.add_movie(path, left, top, width, height)
            except Exception as e:
                print(f"Lỗi khi chèn video: {e}")
        else:
            print(f"Không tìm thấy file video tại: {path}")

    def convert(self, md_text):
        lines = md_text.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # # Title Slide
            if line.startswith('# '):
                self.add_title_slide(line[2:])
            # ## Content Slide
            elif line.startswith('## '):
                self.add_content_slide(line[3:])
            # Bullets
            elif line.startswith(('-', '*')):
                self.add_bullet(line[1:].strip())
            # Images
            elif '![' in line and '](' in line:
                match = re.search(r'!\[.*?\]\((.*?)\)', line)
                if match:
                    self.add_image(match.group(1))
            # Videos (custom syntax)
            elif line.lower().startswith('video:'):
                video_path = line[6:].strip()
                self.add_video(video_path)
            # plain text
            else:
                if self.current_content_placeholder:
                    self.add_bullet(line, level=0)

    def save(self, output_path):
        self.prs.save(output_path)
        print(f"Đã lưu file PowerPoint tại: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 2 or sys.argv[1] != "config_dunp":
        print("Sử dụng: python test/test_gemmar4_ppt.py config_dunp")
        sys.exit(1)

    # Demo Markdown
    demo_md = """
# Giới thiệu về Project Gemma 4
Đây là một dự án demo tạo PowerPoint từ Markdown.

## Các tính năng chính
- Chuyển đổi Markdown thành Slides
- Hỗ trợ ảnh từ URL và local path
- Hỗ trợ nhúng Video local
- Tự động định dạng slide title và nội dung

## Demo Hình ảnh
![Google Logo](https://www.google.com/images/branding/googlelogo/2x/googlelogo_color_272x92dp.png)

## Demo Video
Dưới đây là một slide chứa video (thay đường dẫn thực tế của bạn).
video: test/sample_video.mp4

## Kết luận
Hy vọng công cụ này giúp ích cho việc báo cáo nhanh.
"""
    
    # Tạo folder test nếu chưa có (theo logic user quy định nằm trong folder test)
    if not os.path.exists('test'):
        os.makedirs('test')

    output_file = 'test/output_test_gemmar4_ppt.pptx'
    
    # Mock a video file for testing if it doesn't exist? 
    # (Optional, but let's just run it with what we have)
    
    converter = MarkdownToPPTConverter()
    converter.convert(demo_md)
    converter.save(output_file)
